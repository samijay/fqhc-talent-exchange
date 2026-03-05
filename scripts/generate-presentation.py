#!/usr/bin/env python3
"""
Generate PowerPoint presentation: "The AI-Native Development Workflow"
How one person built a strategic intelligence platform with Claude Code.

Usage: python3 scripts/generate-presentation.py
Output: resources/AI-Native-Development-Workflow.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand Colors ──────────────────────────────────────────────────────
BG_DARK = RGBColor(0x0C, 0x0A, 0x09)       # stone-950
BG_CARD = RGBColor(0x1C, 0x1A, 0x17)       # stone-900
BG_CARD_LIGHT = RGBColor(0x29, 0x25, 0x24) # stone-800
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GRAY = RGBColor(0xA8, 0xA2, 0x9E)      # stone-400
TEXT_DIM = RGBColor(0x78, 0x71, 0x6C)       # stone-500
TEAL = RGBColor(0x2D, 0xD4, 0xBF)          # teal-400
TEAL_DARK = RGBColor(0x0F, 0x76, 0x6E)     # teal-700
AMBER = RGBColor(0xFB, 0xBF, 0x24)         # amber-400
RED = RGBColor(0xF8, 0x71, 0x71)           # red-400
GREEN = RGBColor(0x4A, 0xDE, 0x80)         # green-400
PURPLE = RGBColor(0xC0, 0x84, 0xFC)        # purple-400

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    """Set slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    """Add a text box with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines,
                       font_name="Calibri"):
    """Add a text box with multiple styled lines.
    lines = [(text, font_size, color, bold, alignment), ...]
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, font_size, color, bold, alignment) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)

    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD,
                     border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_stat_card(slide, left, top, value, label, accent=False):
    """Add a stat card with value and label."""
    card_w = Inches(2.5)
    card_h = Inches(1.4)
    add_rounded_rect(slide, left, top, card_w, card_h, BG_CARD)
    add_text_box(slide, left, top + Inches(0.2), card_w, Inches(0.7),
                 value, font_size=36, bold=True,
                 color=AMBER if accent else TEXT_WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.85), card_w, Inches(0.4),
                 label, font_size=12, color=TEXT_DIM,
                 alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Badge
    add_text_box(slide, Inches(4), Inches(0.8), Inches(5.3), Inches(0.5),
                 "⌨  CLAUDE CODE CASE STUDY", font_size=14, color=TEAL,
                 alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.2),
                 "The AI-Native", font_size=54, bold=True,
                 color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10.3), Inches(1.2),
                 "Development Workflow", font_size=54, bold=True,
                 color=TEAL, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(2.5), Inches(3.9), Inches(8.3), Inches(0.8),
                 "How one person built a 110+ feature strategic intelligence\nplatform in 3 weeks using Claude Code",
                 font_size=20, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    # Stat cards
    y = Inches(5.2)
    spacing = Inches(2.8)
    start_x = Inches(1.7)
    add_stat_card(slide, start_x, y, "110+", "Features Shipped", accent=True)
    add_stat_card(slide, start_x + spacing, y, "214", "FQHCs Tracked")
    add_stat_card(slide, start_x + spacing * 2, y, "3", "Weeks", accent=True)
    add_stat_card(slide, start_x + spacing * 3, y, "1", "Developer")


def slide_problem(prs):
    """Slide 2: The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.4),
                 "⚠  THE CRISIS", font_size=14, color=RED)

    add_multiline_text(slide, Inches(1), Inches(1.1), Inches(11), Inches(2), [
        ("California's safety-net workforce", 44, TEXT_WHITE, True, PP_ALIGN.LEFT),
        ("is collapsing", 44, RED, True, PP_ALIGN.LEFT),
    ])

    add_text_box(slide, Inches(1), Inches(3.2), Inches(9), Inches(0.9),
                 "Federally Qualified Health Centers serve 7.8M Californians — mostly uninsured, undocumented, and low-income. They're facing the worst funding crisis in 60 years.",
                 font_size=18, color=TEXT_GRAY)

    # Crisis stat cards
    stats = [
        ("$4.6B", "Federal funding at risk\nfrom H.R. 1"),
        ("3,477+", "Healthcare workers\ndisplaced in CA"),
        ("84%", "Of FQHCs report\nworkforce shortages"),
        ("50%", "Revenue loss projected\nfor some FQHCs"),
    ]
    y = Inches(4.8)
    for i, (val, label) in enumerate(stats):
        x = Inches(1) + Inches(3) * i
        card_w = Inches(2.7)
        card_h = Inches(1.8)
        add_rounded_rect(slide, x, y, card_w, card_h, RGBColor(0x2A, 0x0A, 0x0A),
                         border_color=RGBColor(0x5C, 0x1F, 0x1F))
        add_text_box(slide, x, y + Inches(0.25), card_w, Inches(0.6),
                     val, font_size=32, bold=True, color=RED,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + Inches(0.95), card_w, Inches(0.7),
                     label, font_size=11, color=TEXT_GRAY,
                     alignment=PP_ALIGN.CENTER)


def slide_what_built(prs):
    """Slide 3: What Got Built"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.4),
                 "🚀  THE PRODUCT", font_size=14, color=TEAL)

    add_text_box(slide, Inches(1), Inches(1.1), Inches(11), Inches(1),
                 "A full strategic intelligence platform",
                 font_size=44, bold=True, color=TEXT_WHITE)

    add_text_box(slide, Inches(1), Inches(2.3), Inches(10), Inches(0.6),
                 "Live at fqhctalent.com — built for FQHC executives tracking the political landscape and workforce crisis.",
                 font_size=18, color=TEXT_GRAY)

    features = [
        "214 FQHC Directory", "Intelligence Dashboard", "Layoff Tracker",
        "AI Adoption Tracker", "13 Blog Articles (EN/ES)", "9 Regional Dashboards",
        "220 Strategic Reports", "Resume Builder", "Salary Intelligence",
        "Resilience Scorecard", "Career Assessment", "15 Masterclasses",
        "6 Slash Commands", "Funding Impact Tracker", "Fully Bilingual",
        "FQHC Comparison Tool",
    ]
    cols = 4
    rows = 4
    y_start = Inches(3.3)
    for i, feat in enumerate(features):
        col = i % cols
        row = i // cols
        x = Inches(1) + Inches(2.9) * col
        y = y_start + Inches(0.65) * row
        card_w = Inches(2.7)
        card_h = Inches(0.5)
        add_rounded_rect(slide, x, y, card_w, card_h, BG_CARD)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.05), card_w - Inches(0.3), card_h,
                     f"▸ {feat}", font_size=11, color=TEXT_GRAY)


def slide_reveal(prs):
    """Slide 4: The Reveal"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1.5), Inches(1), Inches(10.3), Inches(1.2),
                 "Built by one person", font_size=52, bold=True,
                 color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2.5), Inches(2.3), Inches(8.3), Inches(0.8),
                 "A self-described beginner with no prior React, Next.js,\nor TypeScript experience.",
                 font_size=20, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    # Code block
    code_w = Inches(7)
    code_h = Inches(2.3)
    code_x = Inches(3.2)
    code_y = Inches(3.5)
    add_rounded_rect(slide, code_x, code_y, code_w, code_h, BG_CARD)

    code_text = '// from CLAUDE.md\n\npreferences:\n  - "I am a beginner — explain things clearly"\n  - "Keep code simple and readable"\n  - "Always test that code compiles"'
    add_text_box(slide, code_x + Inches(0.3), code_y + Inches(0.2),
                 code_w - Inches(0.6), code_h - Inches(0.4),
                 code_text, font_size=14, color=TEAL,
                 font_name="Courier New")

    add_text_box(slide, Inches(1.5), Inches(6.1), Inches(10.3), Inches(0.5),
                 "The secret weapon?  ⌨  Claude Code",
                 font_size=24, bold=True, color=TEXT_WHITE,
                 alignment=PP_ALIGN.CENTER)


def slide_claude_code(prs):
    """Slide 5: What is Claude Code"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.4),
                 "⌨  THE TOOL", font_size=14, color=TEAL)

    add_text_box(slide, Inches(1), Inches(1.1), Inches(11), Inches(1),
                 "Claude Code = AI in your terminal",
                 font_size=44, bold=True, color=TEXT_WHITE)

    add_text_box(slide, Inches(1), Inches(2.3), Inches(10), Inches(0.6),
                 "An agentic coding tool that lives in your terminal. It reads your codebase, writes code, runs commands, and thinks about architecture.",
                 font_size=18, color=TEXT_GRAY)

    # Terminal block (left)
    term_w = Inches(6)
    term_h = Inches(3.5)
    add_rounded_rect(slide, Inches(0.7), Inches(3.5), term_w, term_h, RGBColor(0x09, 0x09, 0x0B),
                     border_color=BG_CARD_LIGHT)
    term_text = '$ claude "add dark mode to settings"\n\nI\'ll plan this implementation first...\n\nReading src/app/settings/page.tsx...\nReading src/lib/theme-config.ts...\n\nPlan:\n1. Add theme toggle component\n2. Wire up next-themes provider\n3. Update CSS variables\n\nShall I proceed? (Y/n)'
    add_text_box(slide, Inches(1), Inches(3.7), term_w - Inches(0.6), term_h - Inches(0.4),
                 term_text, font_size=11, color=GREEN, font_name="Courier New")

    # Capability cards (right)
    caps = [
        ("Reads your entire codebase", "Understands structure, patterns, types", TEAL),
        ("Writes production code", "TypeScript, React, API routes, data files", AMBER),
        ("Runs commands & verifies", "npm run build, git commit, curl APIs", GREEN),
        ("Remembers via CLAUDE.md", "Project memory persists across sessions", PURPLE),
    ]
    for i, (title, desc, color) in enumerate(caps):
        y = Inches(3.5) + Inches(0.85) * i
        x = Inches(7.2)
        card_w = Inches(5.5)
        card_h = Inches(0.75)
        add_rounded_rect(slide, x, y, card_w, card_h, BG_CARD)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.05), card_w - Inches(0.4), Inches(0.35),
                     title, font_size=13, bold=True, color=color)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.38), card_w - Inches(0.4), Inches(0.3),
                     desc, font_size=10, color=TEXT_DIM)


def slide_dev_loop(prs):
    """Slide 6: The Development Loop"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.4),
                 "🔄  THE WORKFLOW", font_size=14, color=TEAL)

    add_text_box(slide, Inches(1), Inches(1.1), Inches(11), Inches(1),
                 "The build loop", font_size=44, bold=True, color=TEXT_WHITE)

    add_text_box(slide, Inches(1), Inches(2.2), Inches(10), Inches(0.6),
                 "Every feature follows the same cycle. Claude Code handles 90% of the execution — I provide direction and domain expertise.",
                 font_size=18, color=TEXT_GRAY)

    steps = [
        ("1. Describe", "Tell Claude what\nto build in plain\nEnglish", TEAL),
        ("2. Plan", "Claude explores\ncodebase, proposes\napproach", RGBColor(0x60, 0xA5, 0xFA)),
        ("3. Build", "Claude writes code\nacross multiple\nfiles", AMBER),
        ("4. Verify", "npm run build —\nmust compile\nclean", GREEN),
        ("5. Ship", "Git commit + push\n→ live on Vercel", PURPLE),
    ]
    y = Inches(3.5)
    for i, (title, desc, color) in enumerate(steps):
        x = Inches(0.8) + Inches(2.5) * i
        card_w = Inches(2.2)
        card_h = Inches(2.2)
        # Colored top accent
        accent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x, y, card_w, Inches(0.08))
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()
        # Card
        add_rounded_rect(slide, x, y + Inches(0.1), card_w, card_h, BG_CARD)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.25), card_w - Inches(0.4), Inches(0.4),
                     title, font_size=16, bold=True, color=color)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.75), card_w - Inches(0.4), Inches(1.2),
                     desc, font_size=12, color=TEXT_GRAY)
        # Arrow between cards
        if i < 4:
            add_text_box(slide, x + card_w, y + Inches(0.9), Inches(0.3), Inches(0.4),
                         "→", font_size=20, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

    # Bottom callout
    add_rounded_rect(slide, Inches(3.5), Inches(6.2), Inches(6.3), Inches(0.7),
                     RGBColor(0x2A, 0x1A, 0x05), border_color=RGBColor(0x5C, 0x3D, 0x0A))
    add_text_box(slide, Inches(3.5), Inches(6.3), Inches(6.3), Inches(0.5),
                 "Average time from idea → deployed feature:  20-45 minutes",
                 font_size=14, color=AMBER, alignment=PP_ALIGN.CENTER)


def slide_slash_commands(prs):
    """Slide 7: Slash Commands"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.4),
                 "⌨  AUTOMATION", font_size=14, color=TEAL)

    add_text_box(slide, Inches(1), Inches(1.1), Inches(11), Inches(1),
                 "Slash commands as infrastructure",
                 font_size=44, bold=True, color=TEXT_WHITE)

    add_text_box(slide, Inches(1), Inches(2.3), Inches(10), Inches(0.6),
                 "Custom commands in .claude/commands/ that encode complex workflows as reusable playbooks.",
                 font_size=18, color=TEXT_GRAY)

    commands = [
        ("/daily-update", "WARN Act + job scan + policy scan + AI + news + regional", "20 searches  •  ~20 min"),
        ("/scan-policy", "Deep legislative scanner: federal, state, local tracking", "10+ sources  •  ~15 min"),
        ("/intel-brief", "Generate weekly newsletter with primary source links", "Aggregation  •  ~10 min"),
        ("/update-layoffs", "Fetch CA WARN Act XLSX, parse, filter, cross-reference", "EDD data  •  ~5 min"),
        ("/scrape-jobs", "Check FQHC career pages via Workday/Lever APIs", "4 APIs  •  ~3 min"),
        ("/draft-blog", "Suggest topics from data, draft bilingual article", "Data + web  •  ~30 min"),
    ]
    cols = 3
    for i, (cmd, desc, meta) in enumerate(commands):
        col = i % cols
        row = i // cols
        x = Inches(0.7) + Inches(4.1) * col
        y = Inches(3.3) + Inches(2) * row
        card_w = Inches(3.8)
        card_h = Inches(1.8)
        add_rounded_rect(slide, x, y, card_w, card_h, BG_CARD)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.15), card_w - Inches(0.4), Inches(0.4),
                     cmd, font_size=16, bold=True, color=TEAL, font_name="Courier New")
        add_text_box(slide, x + Inches(0.2), y + Inches(0.6), card_w - Inches(0.4), Inches(0.7),
                     desc, font_size=11, color=TEXT_GRAY)
        add_text_box(slide, x + Inches(0.2), y + Inches(1.3), card_w - Inches(0.4), Inches(0.3),
                     meta, font_size=10, color=TEXT_DIM)


def slide_daily_pipeline(prs):
    """Slide 8: Daily Pipeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.4),
                 "▶  LIVE DEMO", font_size=14, color=AMBER)

    add_text_box(slide, Inches(1), Inches(1), Inches(11), Inches(0.8),
                 "The daily pipeline", font_size=44, bold=True, color=TEXT_WHITE)

    add_text_box(slide, Inches(1), Inches(1.9), Inches(10), Inches(0.5),
                 "Every morning, one command updates the entire platform.",
                 font_size=18, color=TEXT_GRAY)

    # Terminal
    term_w = Inches(11)
    term_h = Inches(4.5)
    term_x = Inches(1.2)
    term_y = Inches(2.7)

    # Terminal chrome bar
    add_rounded_rect(slide, term_x, term_y, term_w, Inches(0.4), RGBColor(0x1C, 0x1A, 0x17))
    add_text_box(slide, term_x + Inches(0.5), term_y + Inches(0.05), Inches(4), Inches(0.3),
                 "●  ●  ●   claude-code — daily pipeline",
                 font_size=10, color=TEXT_DIM)

    # Terminal body
    add_rounded_rect(slide, term_x, term_y + Inches(0.4), term_w, term_h - Inches(0.4),
                     RGBColor(0x09, 0x09, 0x0B))

    pipeline_text = """$ /daily-update

=== DAILY UPDATE 2026-03-05 ===

Step 1: WARN Act Check
  → Downloading CA EDD WARN XLSX...
  → 126 healthcare entries, 0 FQHC filings

Step 2: Job Scan (4 APIs in parallel)
  → AltaMed: 228  |  FHCSD: 150  |  AHS: 19  |  La Clinica: 179
  → Total: 576 jobs (stable)

Step 3: Legislative & Policy Scan → SD County safety net overhaul (Critical)
Step 3.5: News & Intelligence Scan → +6 new intel items
Step 3.8: Regional Scan (San Diego + Inland Empire) → $300M county funding loss

Step 4.5: Link QC → 4 new URLs verified, 0 broken
Build: PASS ✓    Intel items: 55 total"""

    add_text_box(slide, term_x + Inches(0.3), term_y + Inches(0.55),
                 term_w - Inches(0.6), term_h - Inches(0.7),
                 pipeline_text, font_size=12, color=GREEN, font_name="Courier New")


def slide_conversational(prs):
    """Slide 9: Conversational Development"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.4),
                 "✨  REAL EXAMPLES", font_size=14, color=TEAL)

    add_text_box(slide, Inches(1), Inches(1.1), Inches(11), Inches(1),
                 "Prompts → Features", font_size=44, bold=True, color=TEXT_WHITE)

    add_text_box(slide, Inches(1), Inches(2.1), Inches(10), Inches(0.5),
                 "Real prompts from the build. Each produced a complete, deployed feature.",
                 font_size=18, color=TEXT_GRAY)

    examples = [
        ('"Build a resilience scorecard that scores all 220 FQHCs across 5 dimensions"',
         "→ Scoring engine (350 lines) + interactive page with search, filter, sort, grade distribution viz",
         "~40 min"),
        ('"Add regional intelligence pages — one dashboard per CA region"',
         "→ 9 SSG pages × 2 locales = 18 static pages with per-region stats, resilience, EHR, jobs",
         "~35 min"),
        ('"We missed a massive SF crisis. Build a regional news scanning system."',
         "→ Regional config (9 regions), 5-day rotation, 3 query templates, pipeline integration",
         "~45 min"),
        ('"Make AltaMed incredible"',
         "→ Deep research: Glassdoor 3.6/1.4K, 3K+ staff, 300K patients, Epic, 7 programs, union info",
         "~20 min"),
    ]
    for i, (prompt, result, time) in enumerate(examples):
        y = Inches(2.9) + Inches(1.1) * i
        card_w = Inches(11.3)
        card_h = Inches(0.95)
        add_rounded_rect(slide, Inches(1), y, card_w, card_h, BG_CARD)

        # Number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(1.2), y + Inches(0.2), Inches(0.4), Inches(0.4))
        badge.fill.solid()
        badge.fill.fore_color.rgb = TEAL_DARK
        badge.line.fill.background()
        add_text_box(slide, Inches(1.2), y + Inches(0.2), Inches(0.4), Inches(0.4),
                     str(i + 1), font_size=12, bold=True, color=TEXT_WHITE,
                     alignment=PP_ALIGN.CENTER)

        add_text_box(slide, Inches(1.8), y + Inches(0.1), card_w - Inches(1.5), Inches(0.4),
                     prompt, font_size=12, color=TEAL, font_name="Courier New")
        add_text_box(slide, Inches(1.8), y + Inches(0.5), card_w - Inches(2.5), Inches(0.35),
                     result, font_size=11, color=TEXT_GRAY)
        add_text_box(slide, Inches(11), y + Inches(0.55), Inches(1.2), Inches(0.3),
                     time, font_size=10, color=TEXT_DIM, alignment=PP_ALIGN.RIGHT)


def slide_agents(prs):
    """Slide 10: Parallel Agents"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.4),
                 "🤖  SCALING", font_size=14, color=TEAL)

    add_text_box(slide, Inches(1), Inches(1.1), Inches(11), Inches(1),
                 "Background agents for parallel work",
                 font_size=44, bold=True, color=TEXT_WHITE)

    add_text_box(slide, Inches(1), Inches(2.2), Inches(10), Inches(0.6),
                 "Claude Code can spawn background research agents that work simultaneously. I used this to enrich 220 FQHCs with real data.",
                 font_size=18, color=TEXT_GRAY)

    # Terminal
    term_w = Inches(11)
    term_h = Inches(3.3)
    term_x = Inches(1.2)
    term_y = Inches(3.2)
    add_rounded_rect(slide, term_x, term_y, term_w, Inches(0.4), BG_CARD)
    add_text_box(slide, term_x + Inches(0.5), term_y + Inches(0.05), Inches(5), Inches(0.3),
                 "●  ●  ●   claude-code — parallel enrichment", font_size=10, color=TEXT_DIM)
    add_rounded_rect(slide, term_x, term_y + Inches(0.4), term_w, term_h - Inches(0.4),
                     RGBColor(0x09, 0x09, 0x0B))

    agent_text = """$ Enrich the next batch of HRSA FQHCs

Launching 6 research agents in parallel...

Agent 1: Bay Area FQHCs (14 orgs)      [RUNNING]
Agent 2: Central Valley + SD (10 orgs)  [RUNNING]
Agent 3: Sacramento + Coast (6 orgs)    [RUNNING]
Agent 4: North State + Coast (11 orgs)  [RUNNING]
Agent 5: LA batch 1 (20 orgs)           [RUNNING]
Agent 6: LA batch 2 (20 orgs)           [RUNNING]

→ 6 agents × ~15 min each = 81 FQHCs enriched in parallel"""

    add_text_box(slide, term_x + Inches(0.3), term_y + Inches(0.55),
                 term_w - Inches(0.6), term_h - Inches(0.7),
                 agent_text, font_size=12, color=GREEN, font_name="Courier New")

    # Callout
    add_rounded_rect(slide, Inches(2.5), Inches(6.7), Inches(8.3), Inches(0.6),
                     RGBColor(0x2A, 0x1A, 0x05), border_color=RGBColor(0x5C, 0x3D, 0x0A))
    add_text_box(slide, Inches(2.5), Inches(6.75), Inches(8.3), Inches(0.5),
                 "Without agents: ~15 hours.  With agents: ~15 minutes.",
                 font_size=14, color=AMBER, alignment=PP_ALIGN.CENTER)


def slide_data_moat(prs):
    """Slide 11: The Data Moat"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.4),
                 "🗄  ARCHITECTURE", font_size=14, color=TEAL)

    add_text_box(slide, Inches(1), Inches(1.1), Inches(11), Inches(1),
                 "15+ data files as the intelligence moat",
                 font_size=44, bold=True, color=TEXT_WHITE)

    add_text_box(slide, Inches(1), Inches(2.2), Inches(10), Inches(0.5),
                 "The real product isn't the UI — it's the structured data underneath.",
                 font_size=18, color=TEXT_GRAY)

    files = [
        ("california-fqhcs.ts", "230KB", "220 FQHCs: Glassdoor, programs, EHR, salary, union"),
        ("fqhc-job-listings.ts", "125KB", "156 listings: salary, role, region, requirements"),
        ("funding-impact-data.ts", "123KB", "H.R. 1 timeline, revenue strategies, enrollment"),
        ("fqhc-news-intel.ts", "530 ln", "55 curated intel items, 8 categories, bilingual"),
        ("scope-of-practice.ts", "1684 ln", "10 CA roles, delegation matrix, BPC citations"),
        ("fqhc-movement-history.ts", "1332 ln", "30 events, 5 eras, 8 cross-cultural alliances"),
        ("career-assessment-engine.ts", "40KB+", "55 questions, 5 domains, role scoring"),
        ("fqhc-resilience.ts", "350 ln", "5-dimension scoring for all 220 FQHCs"),
        ("regional-news-sources.ts", "350 ln", "9 regions, outlets, rotation schedule"),
    ]
    cols = 3
    for i, (name, size, desc) in enumerate(files):
        col = i % cols
        row = i // cols
        x = Inches(0.7) + Inches(4.1) * col
        y = Inches(3.1) + Inches(1.45) * row
        card_w = Inches(3.8)
        card_h = Inches(1.3)
        add_rounded_rect(slide, x, y, card_w, card_h, BG_CARD)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.1), card_w - Inches(1.2), Inches(0.3),
                     name, font_size=11, color=TEAL, font_name="Courier New")
        add_text_box(slide, x + card_w - Inches(0.8), y + Inches(0.1), Inches(0.7), Inches(0.3),
                     size, font_size=10, color=TEXT_DIM, alignment=PP_ALIGN.RIGHT)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.5), card_w - Inches(0.3), Inches(0.7),
                     desc, font_size=10, color=TEXT_GRAY)


def slide_claude_md(prs):
    """Slide 12: CLAUDE.md"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.4),
                 "📄  THE SECRET", font_size=14, color=TEAL)

    add_text_box(slide, Inches(1), Inches(1.1), Inches(11), Inches(1),
                 "CLAUDE.md = Project memory",
                 font_size=44, bold=True, color=TEXT_WHITE)

    add_text_box(slide, Inches(1), Inches(2.2), Inches(10), Inches(0.5),
                 "A markdown file checked into your repo that Claude reads at the start of every session.",
                 font_size=18, color=TEXT_GRAY)

    # Left: what's in it
    items = [
        "Project stack & conventions",
        "Every feature ever built (110+ entries)",
        "Data source inventory (15+ files with sizes)",
        "Slash command documentation",
        "Session log with daily summaries",
        "Decisions made (with reasoning)",
        "Nav structure & color palette",
        "Terms glossary (FQHC, ECM, CalAIM...)",
        "Database schema reference",
        "What's NOT built yet (MVP gaps)",
    ]
    add_text_box(slide, Inches(1), Inches(3), Inches(5), Inches(0.4),
                 "What's in ours (900+ lines):", font_size=14, bold=True, color=TEXT_GRAY)
    for i, item in enumerate(items):
        y = Inches(3.5) + Inches(0.35) * i
        add_text_box(slide, Inches(1.3), y, Inches(5), Inches(0.3),
                     f"✓  {item}", font_size=12, color=TEXT_GRAY)

    # Right: session log
    add_rounded_rect(slide, Inches(7), Inches(3), Inches(5.5), Inches(4.2), BG_CARD)
    session_text = """Session Log (excerpt)

2026-03-05  Daily update #11 + regional
            news scanning system
2026-03-04  Locum tenens + SF crisis
            cluster (8 intel items)
2026-03-03  Daily update #8 + branding
2026-03-01  Masterclass + bibliography
2026-02-28  🎉 100th feature shipped!
2026-02-27  Strategic pivot — Rumelt
            framework homepage
...         16 session entries total"""
    add_text_box(slide, Inches(7.3), Inches(3.2), Inches(5), Inches(3.8),
                 session_text, font_size=11, color=TEAL, font_name="Courier New")


def slide_timeline(prs):
    """Slide 13: Timeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.4),
                 "⏱  3 WEEKS", font_size=14, color=AMBER)

    add_text_box(slide, Inches(1), Inches(0.9), Inches(11), Inches(0.8),
                 "The timeline", font_size=44, bold=True, color=TEXT_WHITE)

    # Week 1
    w1_x = Inches(0.5)
    w1_y = Inches(2)
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, w1_x, w1_y, Inches(0.6), Inches(0.6))
    badge.fill.solid()
    badge.fill.fore_color.rgb = TEAL_DARK
    badge.line.fill.background()
    add_text_box(slide, w1_x, w1_y + Inches(0.1), Inches(0.6), Inches(0.4),
                 "W1", font_size=12, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

    w1_items = "Feb 16-18: Foundation\n\n• Career assessment (55 questions, 5 domains)\n• Resume builder (8 templates, PDF)\n• Daily content pipeline (4 slash commands)\n• FQHC directory (90 curated orgs)\n• Security audit (20 vectors, 0 vulns)"
    add_text_box(slide, w1_x + Inches(0.8), w1_y, Inches(3.7), Inches(3),
                 w1_items, font_size=11, color=TEXT_GRAY)

    # Week 2
    w2_x = Inches(5)
    badge2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, w2_x, w1_y, Inches(0.6), Inches(0.6))
    badge2.fill.solid()
    badge2.fill.fore_color.rgb = RGBColor(0xD9, 0x77, 0x06)
    badge2.line.fill.background()
    add_text_box(slide, w2_x, w1_y + Inches(0.1), Inches(0.6), Inches(0.4),
                 "W2", font_size=12, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

    w2_items = "Feb 19-28: Intelligence Pivot\n\n• Directory: 90 → 220 FQHCs (HRSA data)\n• Pivot: job board → executive monitor\n• Rumelt framework homepage\n• 5 strategy pages + 4 interactive vizs\n• Newsletter system (2 tracks)\n• 100th feature shipped 🎉"
    add_text_box(slide, w2_x + Inches(0.8), w1_y, Inches(3.5), Inches(3),
                 w2_items, font_size=11, color=TEXT_GRAY)

    # Week 3
    w3_x = Inches(9.2)
    badge3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, w3_x, w1_y, Inches(0.6), Inches(0.6))
    badge3.fill.solid()
    badge3.fill.fore_color.rgb = RGBColor(0x7C, 0x3A, 0xED)
    badge3.line.fill.background()
    add_text_box(slide, w3_x, w1_y + Inches(0.1), Inches(0.6), Inches(0.4),
                 "W3", font_size=12, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

    w3_items = "Mar 1-5: Scale & Depth\n\n• 15 executive masterclass modules\n• 200+ resource bibliography\n• Regional news scanning (9 regions)\n• SF crisis response: 8 intel items\n• FQHC enrichment: 35+ orgs\n• 220 per-FQHC strategic reports\n• Offboarding toolkit + intake form"
    add_text_box(slide, w3_x + Inches(0.8), w1_y, Inches(3), Inches(3.2),
                 w3_items, font_size=11, color=TEXT_GRAY)

    # Bottom stat
    add_rounded_rect(slide, Inches(4.5), Inches(5.8), Inches(4.3), Inches(1.2),
                     RGBColor(0x0A, 0x2A, 0x25), border_color=TEAL_DARK)
    add_text_box(slide, Inches(4.5), Inches(5.95), Inches(4.3), Inches(0.5),
                 "110+", font_size=36, bold=True, color=TEAL,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(4.5), Inches(6.5), Inches(4.3), Inches(0.4),
                 "features in 16 sessions  •  ~7 per session",
                 font_size=12, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)


def slide_patterns(prs):
    """Slide 14: Key Patterns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.4),
                 "📐  LESSONS", font_size=14, color=TEAL)

    add_text_box(slide, Inches(1), Inches(1.1), Inches(11), Inches(1),
                 "Patterns that worked",
                 font_size=44, bold=True, color=TEXT_WHITE)

    patterns = [
        ("1. Data-first architecture", TEAL,
         "Build the data files first, UI second. The 230KB FQHC directory, the assessment engine, the intel feed — these TypeScript data files ARE the product. Claude Code is excellent at maintaining structured data at scale."),
        ("2. CLAUDE.md as institutional memory", AMBER,
         "Every decision, every feature, every data source — logged in CLAUDE.md. When Claude reads this at session start, it has full context. No re-explaining. No drift. 900+ lines of cumulative knowledge."),
        ("3. Slash commands as playbooks", GREEN,
         "Encode complex multi-step workflows as markdown instructions. The daily-update pipeline is 300+ lines of instructions. Claude follows them like a senior dev following a runbook."),
        ("4. Plan mode for non-trivial work", PURPLE,
         "Always have Claude plan before building anything complex. The planning step catches architectural mistakes, surfaces edge cases, and produces better code on the first pass."),
    ]
    for i, (title, color, desc) in enumerate(patterns):
        col = i % 2
        row = i // 2
        x = Inches(0.7) + Inches(6.2) * col
        y = Inches(2.5) + Inches(2.5) * row
        card_w = Inches(5.8)
        card_h = Inches(2.2)
        add_rounded_rect(slide, x, y, card_w, card_h, BG_CARD)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.2), card_w - Inches(0.6), Inches(0.4),
                     title, font_size=18, bold=True, color=color)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.7), card_w - Inches(0.6), Inches(1.3),
                     desc, font_size=12, color=TEXT_GRAY)


def slide_surprises(prs):
    """Slide 15: What Surprised Me"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.4),
                 "✨  REFLECTIONS", font_size=14, color=AMBER)

    add_text_box(slide, Inches(1), Inches(1.1), Inches(11), Inches(1),
                 "What surprised me",
                 font_size=44, bold=True, color=TEXT_WHITE)

    surprises = [
        ("Speed is non-linear", TEAL,
         "As CLAUDE.md grows and conventions solidify, Claude gets faster, not slower. It remembers your patterns."),
        ("Domain expertise still matters most", AMBER,
         "Claude can write any code — but it can't tell you what to build. Knowing the FQHC industry, understanding the crisis — that's the human edge."),
        ("The product pivoted — and that was fine", GREEN,
         "Started as job board → became strategic intelligence platform. Claude Code handled the architectural shift in one session."),
        ('"Beginner" became irrelevant', PURPLE,
         "Started knowing zero React/Next.js/TypeScript. Three weeks later: 110+ features, SSG, API routes, i18n, data viz. Claude doesn't care about your experience level."),
    ]
    for i, (title, color, desc) in enumerate(surprises):
        y = Inches(2.4) + Inches(1.2) * i
        # Number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(1), y + Inches(0.05), Inches(0.5), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        add_text_box(slide, Inches(1), y + Inches(0.1), Inches(0.5), Inches(0.4),
                     str(i + 1), font_size=14, bold=True, color=TEXT_WHITE,
                     alignment=PP_ALIGN.CENTER)

        add_text_box(slide, Inches(1.8), y, Inches(10), Inches(0.4),
                     title, font_size=18, bold=True, color=TEXT_WHITE)
        add_text_box(slide, Inches(1.8), y + Inches(0.45), Inches(10), Inches(0.6),
                     desc, font_size=13, color=TEXT_GRAY)


def slide_closing(prs):
    """Slide 16: Closing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10.3), Inches(0.9),
                 "One person.", font_size=52, bold=True, color=TEXT_WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(2.1), Inches(10.3), Inches(0.9),
                 "Three weeks.", font_size=52, bold=True, color=TEXT_WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.9),
                 "110+ features.", font_size=52, bold=True, color=AMBER,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2.5), Inches(4.2), Inches(8.3), Inches(0.8),
                 "The tools are here. The question isn't whether AI can\nhelp you build — it's what you'll build with it.",
                 font_size=20, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    # Cards
    card_w = Inches(5)
    card_h = Inches(0.9)
    card_x = Inches(4.2)

    add_rounded_rect(slide, card_x, Inches(5.4), card_w, card_h, BG_CARD)
    add_text_box(slide, card_x, Inches(5.45), card_w, Inches(0.35),
                 "Live site", font_size=12, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, card_x, Inches(5.8), card_w, Inches(0.4),
                 "fqhctalent.com", font_size=20, bold=True, color=TEAL,
                 alignment=PP_ALIGN.CENTER)

    add_rounded_rect(slide, card_x, Inches(6.5), card_w, card_h, BG_CARD)
    add_text_box(slide, card_x, Inches(6.55), card_w, Inches(0.35),
                 "Built with", font_size=12, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, card_x, Inches(6.9), card_w, Inches(0.4),
                 "Claude Code by Anthropic", font_size=20, bold=True, color=TEXT_WHITE,
                 alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_title(prs)          # 1
    slide_problem(prs)        # 2
    slide_what_built(prs)     # 3
    slide_reveal(prs)         # 4
    slide_claude_code(prs)    # 5
    slide_dev_loop(prs)       # 6
    slide_slash_commands(prs)  # 7
    slide_daily_pipeline(prs) # 8
    slide_conversational(prs) # 9
    slide_agents(prs)         # 10
    slide_data_moat(prs)      # 11
    slide_claude_md(prs)      # 12
    slide_timeline(prs)       # 13
    slide_patterns(prs)       # 14
    slide_surprises(prs)      # 15
    slide_closing(prs)        # 16

    output_path = "resources/AI-Native-Development-Workflow.pptx"
    prs.save(output_path)
    print(f"✓ Saved {len(prs.slides)} slides to {output_path}")


if __name__ == "__main__":
    main()
